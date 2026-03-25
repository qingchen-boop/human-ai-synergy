#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""生成人机协作学课程PPT"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# 配色方案
PRIMARY = RGBColor(0x2B, 0x5C, 0x9E)      # 深蓝
SECONDARY = RGBColor(0x4A, 0x90, 0xD9)    # 浅蓝
ACCENT = RGBColor(0xF5, 0xA6, 0x23)       # 橙色
TEXT_COLOR = RGBColor(0x33, 0x33, 0x33)   # 深灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

def add_title_slide(prs, title, subtitle=""):
    """添加标题幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
        tf = sub_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    """添加内容幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部装饰条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        if bullet.startswith("    "):
            p.text = bullet.strip()
            p.level = 1
        else:
            p.text = bullet
            p.level = 0
        
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """添加双栏幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 顶部装饰条
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15)
    )
    header.fill.solid()
    header.fill.fore_color.rgb = PRIMARY
    header.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SECONDARY
    
    # 左栏内容
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = ACCENT
    
    # 右栏内容
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.9), Inches(4.3), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(8)

def add_summary_slide(prs, title, points, footer=""):
    """添加总结幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    background.line.fill.background()
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    
    # 要点
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(4))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "[+] " + point
        p.font.size = Pt(22)
        p.font.color.rgb = TEXT_COLOR
        p.space_after = Pt(16)
    
    if footer:
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.5))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = footer
        p.font.size = Pt(16)
        p.font.color.rgb = SECONDARY
        p.alignment = PP_ALIGN.CENTER

def create_presentation():
    """创建完整演示文稿"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # ===== 第1页：封面 =====
    add_title_slide(prs, "人机协作学", "Personal Empowerment in the AI Era")
    
    # ===== 第2页：课程介绍 =====
    add_content_slide(prs, "为什么学习这门课？", [
        "AI正在改变一切",
        "会用AI ≠ 用好AI",
        "这不只是技术课",
        "学完这门课，你将能够：",
        "    理解AI的能力与局限",
        "    掌握高效人机协作方法",
        "    建立健康的AI使用习惯"
    ])
    
    # ===== 第3页：学习方式 =====
    add_two_column_slide(prs, "我们如何学习？",
        "理论学习", [
            "阅读材料",
            "案例分析",
            "小组讨论"
        ],
        "实践练习", [
            "动手实验",
            "真实任务",
            "反思日志"
        ]
    )
    
    # ===== 第4页：什么是AI？ =====
    add_content_slide(prs, "AI到底是什么？", [
        "定义：人工智能——机器模拟人类智能",
        "分类：",
        "    弱AI（Narrow AI）：专注单一任务",
        "    强AI（General AI）：通用智能（尚未实现）",
        "常见AI类型：机器学习、深度学习、大语言模型",
        "类比：AI就像一个超级学习者",
        "    看过海量例子",
        "    能发现规律",
        "    但不懂道理"
    ])
    
    # ===== 第5页：AI能做什么 =====
    add_two_column_slide(prs, "AI的能力边界",
        "AI擅长的", [
            "数据处理与模式识别",
            "快速检索与信息整合",
            "格式转换与内容改写",
            "模式化任务执行",
            "多语言翻译"
        ],
        "AI不擅长的", [
            "真正的理解与常识",
            "跨领域创新",
            "情感理解与共情",
            "模糊情境判断",
            "价值判断与道德推理"
        ]
    )
    
    # ===== 第6页：AI工作原理 =====
    add_content_slide(prs, "AI是如何工作的？", [
        "以大语言模型为例：输入 -> 处理 -> 输出",
        "Step 1: 理解输入 — 把文字转换成数字（向量）",
        "Step 2: 生成回答 — 基于学习到的语言规律预测下一个词",
        "Step 3: 输出结果 — 把数字转换回文字",
        "类比：就像一个读过图书馆所有书的人",
        "    能模仿各种写作风格",
        "    但并不理解内容的真正含义"
    ])
    
    # ===== 第7页：AI的局限性 =====
    add_content_slide(prs, "AI的致命弱点", [
        "案例1：幻觉问题",
        "    问AI：秦始皇是哪年来中国的？",
        "    AI答：公元前260年 (错误)",
        "案例2：推理错误",
        "    AI可能答错基本的逻辑推理题",
        "核心结论：",
        "    AI很擅长听起来对",
        "    但不保证真的是对",
        "    必须保持批判性思维！"
    ])
    
    # ===== 第8页：模块一小结 =====
    add_summary_slide(prs, "模块一核心收获", [
        "理解AI是什么",
        "知道边界才能用好",
        "保持批判性思维"
    ], "思考：你之前对AI有哪些误解？")
    
    # ===== 第9页：人机协作概念 =====
    add_content_slide(prs, "从使用工具到与AI协作", [
        "演进过程：",
        "    1. 工具时代：人类独立完成任务",
        "    2. 辅助时代：人类使用传统工具",
        "    3. 协作时代：人类+AI共同完成",
        "    4. 增强时代：AI增强人类能力",
        "公式：",
        "    人类独特能力 + AI处理能力 = 最佳结果",
        "    判断力、创造力、共情 + 速度、准确、规模化"
    ])
    
    # ===== 第10页：提示工程基础 =====
    add_content_slide(prs, "与AI沟通的艺术——提示工程", [
        "什么是提示（Prompt）？你对AI说的每一句话",
        "为什么提示很重要？同样的AI，不同的提示，效果天差地别",
        "提示工程三要素：",
        "    1. 角色 — 给AI一个身份",
        "    2. 任务 — 清楚说明要什么",
        "    3. 格式 — 指定输出形式",
        "示例对比：",
        "    弱提示：帮我写一封邮件",
        "    强提示：设定角色、明确任务、指定语气"
    ])
    
    # ===== 第11页：提示工程进阶 =====
    add_content_slide(prs, "进阶技巧", [
        "技巧1：Few-shot 示例",
        "    以下是三封好的道歉邮件示例...按这个风格写",
        "技巧2：思维链",
        "    请一步步思考这个问题",
        "技巧3：迭代优化",
        "    问 -> 看结果 -> 指出问题 -> 追问",
        "技巧4：限制条件",
        "    请用小学生能听懂的话解释"
    ])
    
    # ===== 第12页：任务分解 =====
    add_content_slide(prs, "复杂任务如何处理？", [
        "核心思想：不要让AI一次性完成复杂任务",
        "示例：写一篇市场报告",
        "    Step 1: 让AI收集行业信息",
        "    Step 2: 让AI分析竞争对手",
        "    Step 3: 你确定报告框架",
        "    Step 4: 让AI写初稿",
        "    Step 5: 你修改核心内容",
        "    Step 6: 让AI润色",
        "关键：人类负责判断，AI负责执行"
    ])
    
    # ===== 第13页：反馈与迭代 =====
    add_content_slide(prs, "与AI共创的艺术", [
        "误解：",
        "    AI第一版就是最终版",
        "    一次对话解决所有问题",
        "正确：",
        "    AI生成 -> 你反馈 -> AI调整 -> 循环迭代",
        "迭代沟通示例：",
        "    第一轮：帮我写一个产品介绍",
        "    第二轮：第二段太技术化了，能更口语化一些吗？",
        "    第三轮：很好，但开头不够吸引人...",
        "    ...直到满意"
    ])
    
    # ===== 第14页：模块二小结 =====
    add_summary_slide(prs, "模块二核心要点", [
        "提示工程：角色+任务+格式",
        "任务分解：化繁为简",
        "迭代优化：不追求一次到位",
        "协作思维：AI是伙伴，不是工具"
    ], "实践作业：优化你日常使用AI的对话")
    
    # ===== 第15页：AI偏见 =====
    add_content_slide(prs, "当AI变得不公平", [
        "AI输出的结果，反映的是训练数据中的偏见",
        "案例1：招聘AI的歧视",
        "    AI系统性地降低女性简历的评分",
        "案例2：面部识别的种族问题",
        "    早期系统对深肤色人群识别准确率显著较低",
        "案例3：翻译软件的偏见",
        "    He is a doctor, she is a nurse",
        "核心观点：AI的偏见，本质上是人的偏见的规模化呈现"
    ])
    
    # ===== 第16页：AI与就业 =====
    add_two_column_slide(prs, "AI会抢走你的工作吗？",
        "受冲击职业", [
            "数据录入员",
            "重复性流水线",
            "简单客服",
            "基础翻译"
        ],
        "新兴职业", [
            "AI训练师",
            "人机协作设计师",
            "情感陪伴师",
            "跨文化顾问"
        ]
    )
    
    # ===== 第17页：AI伦理原则 =====
    add_content_slide(prs, "AI时代需要什么原则？", [
        "四大原则：",
        "    1. 有益（Beneficence）— AI应该增进人类福祉",
        "    2. 无害（Non-maleficence）— 避免AI造成伤害",
        "    3. 自主（Autonomy）— 人类保持决策自主权",
        "    4. 公正（Justice）— AI的利益和风险公平分配",
        "讨论问题：在日常生活中，你遇到过AI带来的伦理困境吗？"
    ])
    
    # ===== 第18页：模块三小结 =====
    add_summary_slide(prs, "模块三核心收获", [
        "AI不是完美的，它有偏见、有局限",
        "问题不在AI，在于我们如何使用",
        "意识到AI的偏见，保持批判性思维",
        "支持公正的AI发展，为自己的AI使用负责"
    ])
    
    # ===== 第19页：AI学习 =====
    add_content_slide(prs, "AI作为学习伙伴", [
        "对比：传统学习 vs AI辅助学习",
        "    传统：看书 -> 遇到不懂 -> 查字典/搜索 -> 可能找不到答案",
        "    AI辅助：遇到不懂 -> 随时提问 -> 获得个性化解释",
        "苏格拉底式提问法：",
        "    告诉我答案 -> 错误",
        "    我理解对吗？有什么可能的反驳？ -> 正确",
        "    有没有其他角度？ -> 正确",
        "关键：用AI帮助你思考，不是替你思考"
    ])
    
    # ===== 第20页：AI创作 =====
    add_content_slide(prs, "AI增强你的创造力", [
        "你的独特视角",
        "AI可以帮你做：",
        "    头脑风暴扩展思路",
        "    初稿打破空白页恐惧",
        "    格式转换（口语-书面）",
        "    语言润色",
        "你的独特贡献：",
        "    核心观点、个人经历、真实情感、价值判断",
        "最终作品 = 你的 + AI的"
    ])
    
    # ===== 第21页：AI工作流 =====
    add_content_slide(prs, "把AI融入你的日常工作", [
        "设计步骤：",
        "    Step 1: 盘点高频任务",
        "    Step 2: 识别AI切入点",
        "    Step 3: 设计触发条件",
        "    Step 4: 明确AI角色",
        "    Step 5: 建立反馈循环",
        "案例：内容创作者的AI工作流",
        "    选题 -> AI推荐",
        "    写作 -> AI初稿+你改写",
        "    润色 -> AI语言优化",
        "    发布 -> 你最终审核"
    ])
    
    # ===== 第22页：AI边界 =====
    add_content_slide(prs, "什么时候不该用AI？", [
        "情感支持 — AI没有真正的共情 -> 找朋友/家人/心理咨询",
        "重大人生决策 — AI不了解你的全部 -> 找真人顾问",
        "你想自己享受的过程 — 成长体验很重要 -> 亲力亲为",
        "需要身体力行的事 — 技能需要身体记忆 -> 亲自做",
        "当AI成为逃避 — 困难是成长的一部分 -> 面对挑战",
        "过度依赖AI的信号：",
        "    越来越不想独立思考",
        "    没有AI就无法开始"
    ])
    
    # ===== 第23页：模块四小结 =====
    add_summary_slide(prs, "模块四核心要点", [
        "AI学习：加速器，不是替代品",
        "AI创作：增强剂，核心在你",
        "AI工作流：系统化，更高效",
        "AI边界：知道何时不用"
    ], "作业：设计一个属于你自己的AI工作流")
    
    # ===== 第24页：AI Agent =====
    add_content_slide(prs, "AI Agent——下一波浪潮", [
        "什么是AI Agent？能自主行动、完成多步骤任务的AI系统",
        "对比：",
        "    传统AI：你 -> 发指令 -> AI -> 等待 -> 一步步执行",
        "    Agent AI：你 -> 发任务 -> AI自主规划 -> 执行 -> 反馈 -> 调整",
        "示例场景：",
        "    帮我安排一次公司团建",
        "    -> Agent自动查日历、搜地点、比价格、发邀请",
        "潜在风险：",
        "    自主性越高，控制越难",
        "    错误可能放大",
        "    责任归属模糊"
    ])
    
    # ===== 第25页：前沿与未来 =====
    add_content_slide(prs, "AGI离我们有多远？", [
        "时间线预测（专家分歧极大）：",
        "乐观派（5-10年）：",
        "    Jensen Huang (NVIDIA): 5年内",
        "    很多AI研究员：10-20年内",
        "保守派（50年以上/不可能）：",
        "    Yann LeCun: 至少几十年",
        "    很多科学家：不确定",
        "无论AGI何时到来：",
        "    提前思考AI对人类的影响",
        "    建立正确的AI使用习惯",
        "    保持学习的心态"
    ])
    
    # ===== 第26页：模块五小结 =====
    add_summary_slide(prs, "模块五核心收获", [
        "AI在快速发展，我们无法预测未来",
        "但我们可以：",
        "    保持好奇和学习",
        "    思考AI的伦理和社会影响",
        "    为变化做好准备"
    ])
    
    # ===== 第27页：实践练习 =====
    add_content_slide(prs, "让我们动手做", [
        "练习1：提示工程实战",
        "    用提示工程方法优化一个你常用的AI任务",
        "练习2：人机协作设计",
        "    为你自己的工作场景设计一个AI工作流",
        "练习3：偏见识别",
        "    分析一个AI输出中的潜在偏见",
        "练习4：伦理讨论",
        "    话题：AI生成的内容应该由谁负责？"
    ])
    
    # ===== 第28页：课程总结 =====
    add_summary_slide(prs, "人机协作学——回顾与展望", [
        "人类负责：判断力、创造力、价值观",
        "AI负责：执行力、记忆力、规模化",
        "最终目标：",
        "    成为增强型人类",
        "    而不是AI的附庸"
    ])
    
    # ===== 第29页：结束页 =====
    add_title_slide(prs, "谢谢！", "AI会越来越强，但你的判断力、创造力、同理心\n以及作为人的独特价值，永远不会被取代。\n\n学会与AI协作，成为更好的自己。")
    
    # 保存
    output_path = "/home/xiaoxi/.openclaw/workspace/projects/human-ai-synergy/人机协作学-课件.pptx"
    prs.save(output_path)
    print(f"PPT已生成：{output_path}")
    return output_path

if __name__ == "__main__":
    create_presentation()
